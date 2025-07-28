"""
FastAPI application for VecClean.

Provides HTTP endpoints for document processing with high-performance
text cleaning, deduplication, and vectorization.
"""

from __future__ import annotations

import asyncio
import logging
import time
from pathlib import Path
from typing import Any, Dict, List, Optional

import uvicorn
from fastapi import FastAPI, File, HTTPException, UploadFile, Depends, status
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse

from vecclean import __version__, is_cpp_available, get_version_info
from vecclean.api.schemas import (
    CleanAndEmbedRequest,
    CleanAndEmbedResponse,
    HealthResponse,
    VersionResponse,
    ErrorResponse,
)
from vecclean.api.deps import get_pipeline, get_config
from vecclean.core.config import Config
from vecclean.core.pipeline import Pipeline
from vecclean.core.types import ProcessingStatus, VecCleanError
from vecclean.utils.logging import setup_logging


# Setup logging
logger = logging.getLogger(__name__)

# Configuration constants (loaded from config at startup)
MAX_FILES_PER_REQUEST = 100
MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB per file
MAX_TOTAL_SIZE = 1024 * 1024 * 1024  # 1GB total per request
SUPPORTED_FILE_TYPES = {'.txt', '.pdf', '.docx', '.html', '.htm', '.md', '.pptx'}

# Global state
_connection_pools = {}
_embedding_models = {}
_startup_time = time.time()

# Create FastAPI app
app = FastAPI(
    title="VecClean API",
    description="Ultra-low latency text cleaning, deduplication, and vectorization pipeline",
    version=__version__,
    docs_url="/docs",
    redoc_url="/redoc",
    openapi_url="/openapi.json",
)

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Configure appropriately for production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# Exception handlers
@app.exception_handler(VecCleanError)
async def vecclean_error_handler(request, exc: VecCleanError) -> JSONResponse:
    """Handle VecClean-specific errors."""
    logger.error(f"VecClean error: {exc}")
    return JSONResponse(
        status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
        content=ErrorResponse(
            error="processing_error",
            message=str(exc),
            details={"type": type(exc).__name__}
        ).dict()
    )


@app.exception_handler(Exception)
async def general_error_handler(request, exc: Exception) -> JSONResponse:
    """Handle general exceptions."""
    logger.error(f"Unexpected error: {exc}", exc_info=True)
    return JSONResponse(
        status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
        content=ErrorResponse(
            error="internal_error",
            message="An unexpected error occurred",
            details={"type": type(exc).__name__}
        ).dict()
    )


# Startup and shutdown events
@app.on_event("startup")
async def startup_event() -> None:
    """Initialize application on startup."""
    global _startup_time
    _startup_time = time.time()
    
    logger.info("Starting VecClean API server")
    logger.info(f"Version: {__version__}")
    logger.info(f"C++ backend available: {is_cpp_available()}")
    
    # Initialize dependency injection system
    from vecclean.api.deps import startup_dependencies
    await startup_dependencies()
    
    # Initialize connection pools if configured
    config = await get_config()
    if hasattr(config, 'database') and config.database.enabled:
        from vecclean.api.deps import get_database
        try:
            db_conn = await get_database()
            _connection_pools['database'] = db_conn
            logger.info("Database connection pool initialized")
        except Exception as e:
            logger.warning(f"Failed to initialize database pool: {e}")
    
    # Preload embedding models if configured
    if config.embedding.cache_embeddings:
        try:
            from vecclean.core.embedding import LocalSentenceTransformerEmbedding
            embedding_model = LocalSentenceTransformerEmbedding(
                model_name=config.embedding.model_name,
                device=config.embedding.device,
                cache_embeddings=True
            )
            # Warm up the model with a small text
            await embedding_model.encode(["Warming up embedding model"])
            _embedding_models['default'] = embedding_model
            logger.info(f"Preloaded embedding model: {config.embedding.model_name}")
        except Exception as e:
            logger.warning(f"Failed to preload embedding model: {e}")


@app.on_event("shutdown")
async def shutdown_event() -> None:
    """Clean up resources on shutdown."""
    logger.info("Shutting down VecClean API server")
    
    # Shutdown dependency injection system
    from vecclean.api.deps import shutdown_dependencies
    await shutdown_dependencies()
    
    # Close connection pools
    for pool_name, pool in _connection_pools.items():
        try:
            if hasattr(pool, 'disconnect'):
                await pool.disconnect()
            logger.info(f"Closed {pool_name} connection pool")
        except Exception as e:
            logger.error(f"Error closing {pool_name} pool: {e}")
    
    # Save any pending cache data
    for model_name, model in _embedding_models.items():
        try:
            if hasattr(model, 'save_cache'):
                await model.save_cache()
            logger.info(f"Saved cache for {model_name} embedding model")
        except Exception as e:
            logger.error(f"Error saving cache for {model_name}: {e}")
    
    _connection_pools.clear()
    _embedding_models.clear()


# API Routes

@app.get("/healthz", response_model=HealthResponse)
async def health_check() -> HealthResponse:
    """
    Health check endpoint.
    
    Returns system health status and basic performance metrics.
    """
    try:
        # Basic health checks
        import psutil
        import shutil
        
        checks = {}
        
        # C++ backend check
        checks["cpp_backend"] = "healthy" if is_cpp_available() else "unavailable"
        
        # Memory usage check
        memory = psutil.virtual_memory()
        memory_usage_percent = memory.percent
        if memory_usage_percent < 80:
            checks["memory_usage"] = "normal"
        elif memory_usage_percent < 90:
            checks["memory_usage"] = "high"
        else:
            checks["memory_usage"] = "critical"
        
        # Disk space check
        try:
            disk_usage = shutil.disk_usage("/")
            free_space_percent = (disk_usage.free / disk_usage.total) * 100
            if free_space_percent > 20:
                checks["disk_space"] = "normal"
            elif free_space_percent > 10:
                checks["disk_space"] = "low"
            else:
                checks["disk_space"] = "critical"
        except Exception:
            checks["disk_space"] = "unknown"
        
        # Configuration check
        try:
            config = await get_config()
            checks["configuration"] = "healthy"
        except Exception:
            checks["configuration"] = "error"
        
        # Pipeline check
        try:
            pipeline = await get_pipeline()
            checks["pipeline"] = "healthy"
        except Exception:
            checks["pipeline"] = "error"
        
        # Determine overall status
        overall_status = "healthy"
        if any(status in ["critical", "error"] for status in checks.values()):
            overall_status = "unhealthy"
        elif any(status in ["high", "low", "unavailable"] for status in checks.values()):
            overall_status = "degraded"
        
        return HealthResponse(
            status=overall_status,
            timestamp=int(time.time()),
            version=__version__,
            checks=checks
        )
        
    except Exception as e:
        logger.error(f"Health check failed: {e}")
        return HealthResponse(
            status="unhealthy",
            version=__version__,
            cpp_backend="unknown",
            timestamp=time.time(),
            error=str(e)
        )


@app.get("/version", response_model=VersionResponse)
async def get_version() -> VersionResponse:
    """
    Get version information.
    
    Returns detailed version information including backend availability.
    """
    version_info = get_version_info()
    
    return VersionResponse(
        version=__version__,
        cpp_backend_available=is_cpp_available(),
        python_version=get_version_info()["python_version"],
        build_timestamp=get_version_info().get("build_timestamp", "unknown"),
        git_commit=get_version_info().get("git_commit", "unknown"),
        dependencies=get_version_info().get("dependencies", {})
    )


@app.post("/clean-and-embed", response_model=CleanAndEmbedResponse)
async def clean_and_embed_files(
    files: List[UploadFile] = File(..., description="Files to process"),
    config_override: Optional[str] = None,
    include_embeddings: bool = True,
    include_metadata: bool = True,
    output_format: str = "json",
    pipeline: Pipeline = Depends(get_pipeline)
) -> CleanAndEmbedResponse:
    """
    Clean and embed uploaded files.
    
    Main processing endpoint that accepts multiple files and returns
    cleaned, deduplicated, and vectorized chunks.
    
    Args:
        files: List of files to process
        config_override: Optional JSON string to override configuration
        include_embeddings: Whether to include embedding vectors in response
        include_metadata: Whether to include metadata in response  
        output_format: Output format preference
        pipeline: Processing pipeline (injected dependency)
        
    Returns:
        Processing results with cleaned chunks and statistics
        
    Raises:
        HTTPException: If processing fails or files are invalid
    """
    start_time = time.time()
    
    try:
        # Validate inputs
        if not files:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="No files provided"
            )
        
        # File validation
        if len(files) > MAX_FILES_PER_REQUEST:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Too many files. Maximum {MAX_FILES_PER_REQUEST} files allowed per request."
            )
        
        # Check total size
        total_size = sum(len(await file.read()) for file in files)
        for file in files:
            await file.seek(0)  # Reset file pointers
        
        if total_size > MAX_TOTAL_SIZE:
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"Total file size too large. Maximum {MAX_TOTAL_SIZE // (1024*1024)}MB allowed."
            )
        
        # Validate individual file sizes
        for file in files:
            if file.size and file.size > MAX_FILE_SIZE:
                raise HTTPException(
                    status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                    detail=f"File '{file.filename}' is too large. Maximum {MAX_FILE_SIZE // (1024*1024)}MB allowed per file."
                )
            
            # Validate file type
            if file.filename:
                file_ext = Path(file.filename).suffix.lower()
                if file_ext not in SUPPORTED_FILE_TYPES:
                    raise HTTPException(
                        status_code=status.HTTP_415_UNSUPPORTED_MEDIA_TYPE,
                        detail=f"Unsupported file type '{file_ext}'. Supported types: {', '.join(SUPPORTED_FILE_TYPES)}"
                    )
        
        logger.info(f"Processing {len(files)} files: {[f.filename for f in files]}")
        
        # Create temporary directory for file processing
        import tempfile
        from pathlib import Path
        
        temp_dir = Path(tempfile.mkdtemp(prefix="vecclean_"))
        
        try:
            # Process files based on type with proper temporary storage
            all_chunks = []
            file_results = []
            
            for i, file in enumerate(files):
                try:
                    # Save uploaded file to temporary location
                    file_extension = Path(file.filename).suffix.lower() if file.filename else '.txt'
                    temp_file_path = temp_dir / f"file_{i}{file_extension}"
                    
                    # Write file content to temporary location
                    content = await file.read()
                    with open(temp_file_path, 'wb') as temp_file:
                        temp_file.write(content)
                    
                    # Process based on file type
                    if file_extension in ['.txt', '.md']:
                        result = await _process_text_file(temp_file_path, pipeline)
                    elif file_extension in ['.pdf']:
                        result = await _process_pdf_file(temp_file_path, pipeline)
                    elif file_extension in ['.docx']:
                        result = await _process_docx_file(temp_file_path, pipeline)
                    elif file_extension in ['.html', '.htm']:
                        result = await _process_html_file(temp_file_path, pipeline)
                    elif file_extension in ['.pptx']:
                        result = await _process_pptx_file(temp_file_path, pipeline)
                    else:
                        # Fallback to text processing
                        result = await _process_text_file(temp_file_path, pipeline)
                    
                    if result and result.chunks:
                        all_chunks.extend(result.chunks)
                        file_results.append({
                            "filename": file.filename,
                            "status": "success",
                            "chunks": len(result.chunks)
                        })
                    else:
                        file_results.append({
                            "filename": file.filename,
                            "status": "failed",
                            "error": "No content extracted"
                        })
                        
                except Exception as e:
                    logger.error(f"Failed to process file {file.filename}: {e}")
                    file_results.append({
                        "filename": file.filename,
                        "status": "failed",
                        "error": str(e)
                    })
                    continue
                finally:
                    # Clean up temporary file
                    if temp_file_path.exists():
                        temp_file_path.unlink()
        
        finally:
            # Clean up temporary directory
            import shutil
            shutil.rmtree(temp_dir, ignore_errors=True)
        
        # Create combined result
        from vecclean.core.types import ProcessingResult, ProcessingStats
        result = ProcessingResult(
            status=ProcessingStatus.COMPLETED if all_chunks else ProcessingStatus.FAILED,
            chunks=all_chunks,
            stats=ProcessingStats(
                total_files=len(files),
                successful_files=len([f for f in files if any(c.source_document and c.source_document.filename == f.filename for c in all_chunks)]),
                failed_files=len(files) - len(all_chunks),
                total_chunks=len(all_chunks),
                total_processing_time=time.time() - start_time
            ),
            processing_timestamp=time.time()
        )
        
        # Format response
        processing_time = time.time() - start_time
        
        response = CleanAndEmbedResponse(
            status=result.status.value,
            chunks=[
                chunk.to_dict(include_embedding=include_embeddings)
                for chunk in result.chunks
            ],
            statistics=result.stats.to_dict() if include_metadata else None,
            metadata={
                "processing_time_seconds": processing_time,
                "files_processed": len(files),
                "backend_used": "cpp" if is_cpp_available() else "python",
                "config_version": config.version if hasattr(config, 'version') else "1.0.0",
                "chunking_strategy": config.chunking.strategy,
                "max_chunk_size": config.chunking.max_chunk_size,
                "embedding_model": config.embedding.model_name
            } if include_metadata else None,
            warnings=result.warnings,
            errors=result.errors
        )
        
        logger.info(
            f"Successfully processed {len(files)} files in {processing_time:.3f}s, "
            f"generated {len(result.chunks)} chunks"
        )
        
        return response
        
    except HTTPException:
        # Re-raise HTTP exceptions as-is
        raise
    except Exception as e:
        logger.error(f"Processing failed: {e}", exc_info=True)
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Processing failed: {str(e)}"
        )


@app.post("/clean-text", response_model=Dict[str, Any])
async def clean_text_only(
    text: str,
    config_override: Optional[str] = None,
    pipeline: Pipeline = Depends(get_pipeline)
) -> Dict[str, Any]:
    """
    Clean text without chunking or embedding.
    
    Lightweight endpoint for text cleaning only.
    """
    try:
        # Use the normalizer directly for faster processing
        normalizer = pipeline.get_normalizer()
        cleaned_text = await normalizer.clean_text(text)
        
        return {
            "original_text": text,
            "cleaned_text": cleaned_text,
            "statistics": {
                "original_length": len(text),
                "cleaned_length": len(cleaned_text),
                "compression_ratio": len(cleaned_text) / len(text) if len(text) > 0 else 0,
                "backend_used": "cpp" if normalizer.use_cpp else "python"
            }
        }
        
    except Exception as e:
        logger.error(f"Text cleaning failed: {e}")
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Text cleaning failed: {str(e)}"
        )


# Development and debug endpoints (disable in production)

@app.get("/debug/config")
async def get_debug_config(config: Config = Depends(get_config)) -> Dict[str, Any]:
    """Get current configuration (debug only)."""
    return config.to_dict()


@app.get("/debug/capabilities")
async def get_debug_capabilities() -> Dict[str, Any]:
    """Get system capabilities (debug only)."""
    capabilities = {
        "version_info": get_version_info(),
        "cpp_available": is_cpp_available(),
    }
    
    # Add C++ capabilities if available
    if is_cpp_available():
        try:
            import vecclean_cpp
            capabilities["cpp_capabilities"] = vecclean_cpp.get_capabilities()
        except Exception as e:
            capabilities["cpp_error"] = str(e)
    
    return capabilities


# File processing helper functions
async def _process_text_file(file_path: Path, pipeline: Pipeline) -> Optional[Any]:
    """Process text files (.txt, .md)."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        
        return await pipeline.process_text(
            text=text,
            filename=file_path.name,
        )
    except Exception as e:
        logger.error(f"Error processing text file {file_path}: {e}")
        return None


async def _process_pdf_file(file_path: Path, pipeline: Pipeline) -> Optional[Any]:
    """Process PDF files."""
    try:
        # Try PyPDF2 first
        import PyPDF2
        
        text = ""
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        
        if not text.strip():
            # Try pdfplumber as fallback
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
            except ImportError:
                logger.warning("pdfplumber not available for enhanced PDF processing")
        
        return await pipeline.process_text(
            text=text,
            filename=file_path.name,
        )
    except Exception as e:
        logger.error(f"Error processing PDF file {file_path}: {e}")
        return None


async def _process_docx_file(file_path: Path, pipeline: Pipeline) -> Optional[Any]:
    """Process DOCX files."""
    try:
        from docx import Document
        
        doc = Document(file_path)
        text = ""
        
        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        
        return await pipeline.process_text(
            text=text,
            filename=file_path.name,
        )
    except Exception as e:
        logger.error(f"Error processing DOCX file {file_path}: {e}")
        return None


async def _process_html_file(file_path: Path, pipeline: Pipeline) -> Optional[Any]:
    """Process HTML files."""
    try:
        from bs4 import BeautifulSoup
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Extract text
        text = soup.get_text()
        
        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        text = '\n'.join(chunk for chunk in chunks if chunk)
        
        return await pipeline.process_text(
            text=text,
            filename=file_path.name,
        )
    except Exception as e:
        logger.error(f"Error processing HTML file {file_path}: {e}")
        return None


async def _process_pptx_file(file_path: Path, pipeline: Pipeline) -> Optional[Any]:
    """Process PPTX files."""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return await pipeline.process_text(
            text=text,
            filename=file_path.name,
        )
    except Exception as e:
        logger.error(f"Error processing PPTX file {file_path}: {e}")
        return None


# Configuration loading helper
async def _load_api_config():
    """Load API configuration from main config."""
    config = await get_config()
    
    global MAX_FILES_PER_REQUEST, MAX_FILE_SIZE, MAX_TOTAL_SIZE
    
    # Override defaults with config values if available
    if hasattr(config, 'api'):
        MAX_FILES_PER_REQUEST = getattr(config.api, 'max_files_per_request', MAX_FILES_PER_REQUEST)
        MAX_FILE_SIZE = getattr(config.api, 'max_file_size', MAX_FILE_SIZE)
        MAX_TOTAL_SIZE = getattr(config.api, 'max_total_size', MAX_TOTAL_SIZE)


if __name__ == "__main__":
    # Set start time for uptime calculation
    app.state.start_time = time.time()
    
    # Setup logging
    setup_logging()
    
    # Run the server
    uvicorn.run(
        "vecclean.api.main:app",
        host="0.0.0.0",
        port=8000,
        reload=True,
        log_level="info"
    )

# Future enhancements to be implemented in v2:
# - Authentication and authorization (JWT, OAuth)
# - Advanced rate limiting with Redis
# - Request/response compression (gzip, brotli)
# - Streaming responses for large file processing
# - Comprehensive metrics and monitoring
# - Async background processing with queues
# - Webhook support for processing notifications  
# - API versioning (v1, v2)
# - Enhanced OpenAPI schema customization
# - Intelligent request caching with TTL 