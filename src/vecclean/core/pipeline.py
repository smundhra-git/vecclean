"""
Main processing pipeline for VecClean.

This module orchestrates the complete document processing workflow:
1. File ingestion and text extraction
2. Text cleaning and normalization
3. Deduplication at sentence and chunk levels
4. Intelligent text chunking with overlap
5. Embedding generation with caching
6. Output formatting and writing

The pipeline is designed for maximum performance with async I/O for Python
operations and delegation to C++ for CPU-intensive tasks.
"""

from __future__ import annotations

import asyncio
import logging
import time
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

from vecclean.core.config import Config
from vecclean.core.embedding import EmbeddingModel, LocalSentenceTransformerEmbedding
from vecclean.core.types import (
    CleanedChunk,
    DocumentMetadata,
    FileType,
    ProcessingResult,
    ProcessingStats,
    ProcessingStatus,
    VecCleanError,
    ProcessingError,
    FilePath,
    JsonDict,
)
from vecclean.utils.hashing import generate_content_hash
from vecclean.utils.timing import Timer

# Import processing modules (will be implemented as stubs)
from vecclean.ingest.detect import FileTypeDetector
from vecclean.clean.normalize import TextNormalizer
from vecclean.dedup.sentence_dedup import SentenceDeduplicator
from vecclean.dedup.chunk_dedup import ChunkDeduplicator
from vecclean.chunk.splitter import TextChunker


logger = logging.getLogger(__name__)


class Pipeline:
    """
    Main processing pipeline for VecClean.
    
    Coordinates all processing steps from raw files to vectorized chunks.
    Designed for high throughput with async operations and C++ acceleration.
    """
    
    def __init__(self, config: Optional[Config] = None) -> None:
        """
        Initialize the processing pipeline.
        
        Args:
            config: Processing configuration. If None, loads default config.
        """
        # Load configuration
        self.config = config or Config()
        
        # Initialize timing and statistics
        self._timer = Timer()
        self._stats = ProcessingStats()
        
        # Initialize processing components
        self._detector = FileTypeDetector()
        self._normalizer = TextNormalizer(self.config.cleaning)
        self._sentence_dedup = SentenceDeduplicator(self.config.dedup)
        self._chunk_dedup = ChunkDeduplicator(self.config.dedup)
        self._chunker = TextChunker(self.config.chunking)
        
        # Initialize embedding model
        self._embedding_model: Optional[EmbeddingModel] = None
        
        # Content cache for deduplication
        self._content_cache: Dict[str, CleanedChunk] = {}
        
        logger.info(f"Pipeline initialized with config: {self.config.to_dict()}")
    
    async def process_files(
        self, 
        file_paths: List[FilePath],
        output_path: Optional[FilePath] = None
    ) -> ProcessingResult:
        """
        Process multiple files through the complete pipeline.
        
        Args:
            file_paths: List of file paths to process
            output_path: Optional output path for results
            
        Returns:
            Processing results with chunks, stats, and metadata
            
        Raises:
            ProcessingError: If processing fails
        """
        logger.info(f"Starting pipeline processing for {len(file_paths)} files")
        
        # Initialize statistics
        self._stats = ProcessingStats()
        self._stats.total_files = len(file_paths)
        
        # Start overall timing
        with self._timer.time_operation("total_processing"):
            try:
                # Process files in batches for memory efficiency
                all_chunks: List[CleanedChunk] = []
                batch_size = self.config.processing.batch_size
                
                for i in range(0, len(file_paths), batch_size):
                    batch_files = file_paths[i:i + batch_size]
                    batch_chunks = await self._process_file_batch(batch_files)
                    all_chunks.extend(batch_chunks)
                    
                    # Log progress
                    processed_count = min(i + batch_size, len(file_paths))
                    logger.info(f"Processed {processed_count}/{len(file_paths)} files")
                
                # Final deduplication across all chunks
                if self.config.dedup.chunk_dedup and len(all_chunks) > 1:
                    with self._timer.time_operation("final_deduplication"):
                        all_chunks = await self._deduplicate_final_chunks(all_chunks)
                
                # Generate final statistics
                self._finalize_stats(all_chunks)
                
                # Create result
                result = ProcessingResult(
                    status=ProcessingStatus.COMPLETED,
                    chunks=all_chunks,
                    stats=self._stats,
                    processing_timestamp=time.time(),
                    config_used=self.config.to_dict()
                )
                
                # Write output if requested
                if output_path:
                    await self._write_output(result, output_path)
                
                logger.info(f"Pipeline completed successfully. Processed {len(all_chunks)} chunks")
                return result
                
            except Exception as e:
                logger.error(f"Pipeline processing failed: {e}")
                self._stats.failed_files = self._stats.total_files - self._stats.successful_files
                
                return ProcessingResult(
                    status=ProcessingStatus.FAILED,
                    chunks=[],
                    stats=self._stats,
                    errors=[str(e)],
                    processing_timestamp=time.time(),
                    config_used=self.config.to_dict()
                )
    
    async def process_single_file(
        self, 
        file_path: FilePath,
        file_content: Optional[bytes] = None
    ) -> ProcessingResult:
        """
        Process a single file through the pipeline.
        
        Args:
            file_path: Path to the file
            file_content: Optional file content (if already loaded)
            
        Returns:
            Processing results
        """
        return await self.process_files([file_path])
    
    async def process_text(
        self,
        text: str,
        filename: str = "text_input",
        file_type: FileType = FileType.TXT
    ) -> ProcessingResult:
        """
        Process raw text through the pipeline.
        
        Args:
            text: Raw text to process
            filename: Virtual filename for metadata
            file_type: File type for processing hints
            
        Returns:
            Processing results
        """
        logger.info(f"Processing raw text: {len(text)} characters")
        
        with self._timer.time_operation("text_processing"):
            try:
                # Create metadata
                metadata = DocumentMetadata(
                    filename=filename,
                    file_type=file_type,
                    file_size=len(text.encode('utf-8')),
                    word_count=len(text.split()),
                    processing_timestamp=time.time()
                )
                
                # Process through pipeline steps
                chunks = await self._process_extracted_text(text, metadata)
                
                # Generate statistics
                stats = ProcessingStats(
                    total_files=1,
                    successful_files=1,
                    total_chunks=len(chunks),
                    total_text_length=len(text),
                    total_word_count=len(text.split()),
                    total_processing_time=self._timer.get_total_time()
                )
                
                return ProcessingResult(
                    status=ProcessingStatus.COMPLETED,
                    chunks=chunks,
                    stats=stats,
                    processing_timestamp=time.time(),
                    config_used=self.config.to_dict()
                )
                
            except Exception as e:
                logger.error(f"Text processing failed: {e}")
                return ProcessingResult(
                    status=ProcessingStatus.FAILED,
                    chunks=[],
                    stats=ProcessingStats(total_files=1, failed_files=1),
                    errors=[str(e)],
                    processing_timestamp=time.time()
                )
    
    async def _process_file_batch(self, file_paths: List[FilePath]) -> List[CleanedChunk]:
        """Process a batch of files concurrently."""
        tasks = [self._process_single_file_internal(fp) for fp in file_paths]
        batch_results = await asyncio.gather(*tasks, return_exceptions=True)
        
        all_chunks = []
        for result in batch_results:
            if isinstance(result, Exception):
                logger.error(f"File processing error: {result}")
                self._stats.failed_files += 1
            else:
                all_chunks.extend(result)
                self._stats.successful_files += 1
        
        return all_chunks
    
    async def _process_single_file_internal(self, file_path: FilePath) -> List[CleanedChunk]:
        """Process a single file internally."""
        file_path = Path(file_path)
        
        # Implement file size check
        max_size = self.config.processing.max_file_size_mb * 1024 * 1024
        file_size = file_path.stat().st_size
        if file_size > max_size:
            raise ProcessingError(f"File too large: {file_size / 1024 / 1024:.1f}MB > {max_size / 1024 / 1024}MB")
        
        # Implement timeout handling
        timeout = self.config.processing.timeout_seconds
        
        try:
            async with asyncio.timeout(timeout):
                # Detect file type
                file_type = await self._detector.detect_file_type(file_path)
                
                # Extract text and metadata
                text, metadata = await self._extract_text_from_file(file_path, file_type)
                
                # Process extracted text
                chunks = await self._process_extracted_text(text, metadata)
                
                self._stats.total_size_bytes += file_size
                return chunks
        except asyncio.TimeoutError:
            raise ProcessingError(f"Processing timeout ({timeout}s) exceeded for file: {file_path}")
        except Exception as e:
            logger.error(f"Failed to process file {file_path}: {e}")
            raise ProcessingError(f"File processing failed: {e}") from e
    
    async def _extract_text_from_file(
        self, 
        file_path: Path, 
        file_type: FileType
    ) -> tuple[str, DocumentMetadata]:
        """Extract text and metadata from different file formats."""
        
        try:
            if file_type == FileType.PDF:
                text = await self._extract_pdf_text(file_path)
            elif file_type == FileType.DOCX:
                text = await self._extract_docx_text(file_path)
            elif file_type == FileType.HTML:
                text = await self._extract_html_text(file_path)
            elif file_type == FileType.PPTX:
                text = await self._extract_pptx_text(file_path)
            elif file_type == FileType.EMAIL:
                text = await self._extract_email_text(file_path)
            else:
                # Default to plain text
                text = await self._extract_plain_text(file_path)
            
            metadata = DocumentMetadata(
                filename=file_path.name,
                file_type=file_type,
                file_size=file_path.stat().st_size,
                word_count=len(text.split()),
                processing_timestamp=time.time()
            )
            
            return text, metadata
            
        except Exception as e:
            logger.error(f"Failed to extract text from {file_path}: {e}")
            raise ProcessingError(f"Text extraction failed: {e}") from e
    
    async def _extract_pdf_text(self, file_path: Path) -> str:
        """Extract text from PDF files."""
        try:
            import PyPDF2
            
            text = ""
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            
            # Fallback to pdfplumber if no text extracted
            if not text.strip():
                try:
                    import pdfplumber
                    with pdfplumber.open(file_path) as pdf:
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + "\n"
                except ImportError:
                    logger.warning("pdfplumber not available for enhanced PDF processing")
            
            return text
        except Exception as e:
            raise ProcessingError(f"PDF extraction failed: {e}")
    
    async def _extract_docx_text(self, file_path: Path) -> str:
        """Extract text from DOCX files."""
        try:
            from docx import Document
            
            doc = Document(file_path)
            text = ""
            
            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
            
            return text
        except Exception as e:
            raise ProcessingError(f"DOCX extraction failed: {e}")
    
    async def _extract_html_text(self, file_path: Path) -> str:
        """Extract text from HTML files."""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Remove script and style elements
            for script in soup(["script", "style", "nav", "footer"]):
                script.decompose()
            
            # Extract text
            text = soup.get_text()
            
            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)
            
            return text
        except Exception as e:
            raise ProcessingError(f"HTML extraction failed: {e}")
    
    async def _extract_pptx_text(self, file_path: Path) -> str:
        """Extract text from PPTX files."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            return text
        except Exception as e:
            raise ProcessingError(f"PPTX extraction failed: {e}")
    
    async def _extract_email_text(self, file_path: Path) -> str:
        """Extract text from email files."""
        try:
            import email
            
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                msg = email.message_from_file(f)
            
            text = ""
            
            # Extract subject
            if msg['Subject']:
                text += f"Subject: {msg['Subject']}\n\n"
            
            # Extract body
            if msg.is_multipart():
                for part in msg.walk():
                    if part.get_content_type() == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            text += payload.decode('utf-8', errors='ignore')
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    text += payload.decode('utf-8', errors='ignore')
            
            return text
        except Exception as e:
            raise ProcessingError(f"Email extraction failed: {e}")
    
    async def _extract_plain_text(self, file_path: Path) -> str:
        """Extract text from plain text files."""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            raise ProcessingError(f"Plain text extraction failed: {e}")
    
    async def _process_extracted_text(
        self, 
        text: str, 
        metadata: DocumentMetadata
    ) -> List[CleanedChunk]:
        """Process extracted text through the cleaning and chunking pipeline."""
        
        # Step 1: Text cleaning and normalization
        with self._timer.time_operation("text_cleaning"):
            cleaned_text = await self._normalizer.clean_text(text)
        print("[DEBUG] Cleaned text:")
        print(cleaned_text)
        
        # Step 2: Sentence-level deduplication
        if self.config.dedup.sentence_dedup:
            with self._timer.time_operation("sentence_deduplication"):
                cleaned_text = await self._sentence_dedup.deduplicate_sentences(cleaned_text)
        print("[DEBUG] After sentence deduplication:")
        print(cleaned_text)
        
        # Step 3: Text chunking
        with self._timer.time_operation("chunking"):
            chunk_data = await self._chunker.chunk_text(
                cleaned_text,
                chunk_size=self.config.chunking.chunk_size,
                overlap=self.config.chunking.chunk_overlap
            )
        print(f"[DEBUG] Number of chunks: {len(chunk_data)}")
        for i, (chunk_text, start_char, end_char) in enumerate(chunk_data):
            print(f"  Chunk {i}: {chunk_text}")
        
        # Step 4: Create chunk objects
        chunks = []
        for i, (chunk_text, start_char, end_char) in enumerate(chunk_data):
            chunk_hash = generate_content_hash(chunk_text)
            
            # Check cache for duplicates
            if chunk_hash in self._content_cache:
                logger.debug(f"Found duplicate chunk: {chunk_hash}")
                self._stats.duplicate_chunks += 1
                continue
            
            chunk = CleanedChunk(
                chunk_id=f"{metadata.filename}_{i}",
                text=chunk_text,
                text_hash=chunk_hash,
                chunk_index=i,
                start_char=start_char,
                end_char=end_char,
                source_document=metadata
            )
            
            chunks.append(chunk)
            self._content_cache[chunk_hash] = chunk
        
        # Step 5: Generate embeddings
        if chunks:
            await self._generate_embeddings(chunks)
        
        return chunks
    
    async def _generate_embeddings(self, chunks: List[CleanedChunk]) -> None:
        """Generate embeddings for chunks."""
        # Skip embedding generation if model_name is "none"
        if self.config.embedding.model_name == "none":
            logger.info("Embedding generation disabled (model_name='none')")
            return
            
        if not self._embedding_model:
            self._embedding_model = LocalSentenceTransformerEmbedding(
                model_name=self.config.embedding.model_name,
                device=self.config.embedding.device,
                cache_embeddings=self.config.embedding.cache_embeddings
            )
        
        with self._timer.time_operation("embedding_generation"):
            texts = [chunk.text for chunk in chunks]
            embeddings = await self._embedding_model.encode(texts)
            
            for chunk, embedding in zip(chunks, embeddings):
                chunk.embedding = embedding
                chunk.embedding_model = self._embedding_model.get_model_name()
    
    async def _deduplicate_final_chunks(
        self, 
        chunks: List[CleanedChunk]
    ) -> List[CleanedChunk]:
        """Perform final cross-chunk deduplication using ChunkDeduplicator."""
        if not chunks or len(chunks) < 2:
            return chunks
        
        logger.info(f"Performing final deduplication on {len(chunks)} chunks")
        
        try:
            # Use the chunk deduplicator for cross-document deduplication
            deduplicated_chunks = await self._chunk_dedup.deduplicate_chunks(chunks)
            
            removed_count = len(chunks) - len(deduplicated_chunks)
            if removed_count > 0:
                logger.info(f"Final deduplication removed {removed_count} duplicate chunks")
                self._stats.duplicate_chunks += removed_count
            
            return deduplicated_chunks
            
        except Exception as e:
            logger.error(f"Final deduplication failed: {e}")
            return chunks
    
    def _finalize_stats(self, chunks: List[CleanedChunk]) -> None:
        """Calculate final processing statistics."""
        self._stats.total_chunks = len(chunks)
        self._stats.total_text_length = sum(len(chunk.text) for chunk in chunks)
        self._stats.total_word_count = sum(chunk.word_count for chunk in chunks)
        self._stats.total_processing_time = self._timer.get_total_time()
        
        if self._stats.total_files > 0:
            self._stats.avg_processing_time_per_file = (
                self._stats.total_processing_time / self._stats.total_files
            )
        
        if self._stats.total_size_bytes > 0:
            self._stats.avg_processing_time_per_mb = (
                self._stats.total_processing_time / (self._stats.total_size_bytes / 1024 / 1024)
            )
        
        # Calculate compression ratio
        original_length = sum(
            chunk.source_document.word_count or 0 
            for chunk in chunks 
            if chunk.source_document
        )
        if original_length > 0:
            self._stats.compression_ratio = self._stats.total_word_count / original_length
    
    async def _write_output(self, result: ProcessingResult, output_path: FilePath) -> None:
        """Write processing results to output file in various formats."""
        output_path = Path(output_path)
        output_format = self.config.output.format
        
        try:
            if output_format == "jsonl":
                await self._write_jsonl_output(result, output_path)
            elif output_format == "parquet":
                await self._write_parquet_output(result, output_path)
            elif output_format == "json":
                await self._write_json_output(result, output_path)
            else:
                logger.warning(f"Unknown output format: {output_format}, defaulting to JSONL")
                await self._write_jsonl_output(result, output_path)
                
            logger.info(f"Successfully wrote {len(result.chunks)} chunks to {output_path}")
            
        except Exception as e:
            logger.error(f"Failed to write output to {output_path}: {e}")
            raise ProcessingError(f"Output writing failed: {e}")
    
    async def _write_jsonl_output(self, result: ProcessingResult, output_path: Path) -> None:
        """Write results in JSONL format."""
        import json
        import aiofiles
        
        async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
            for chunk in result.chunks:
                chunk_data = {
                    "chunk_id": chunk.chunk_id,
                    "text": chunk.text if self.config.output.include_text else None,
                    "text_hash": chunk.text_hash,
                    "chunk_index": chunk.chunk_index,
                    "start_char": chunk.start_char,
                    "end_char": chunk.end_char,
                    "char_count": chunk.char_count,
                    "word_count": chunk.word_count,
                    "token_count": chunk.token_count,
                    "embedding": chunk.embedding if self.config.output.include_embeddings else None,
                    "embedding_model": chunk.embedding_model,
                    "source_document": chunk.source_document.to_dict() if chunk.source_document else None
                }
                
                # Remove None values if not including metadata
                if not self.config.output.include_metadata:
                    chunk_data = {k: v for k, v in chunk_data.items() if v is not None}
                
                await f.write(json.dumps(chunk_data) + '\n')
    
    async def _write_json_output(self, result: ProcessingResult, output_path: Path) -> None:
        """Write results in JSON format."""
        import json
        import aiofiles
        
        output_data = {
            "status": result.status.value,
            "chunks": [
                chunk.to_dict(
                    include_text=self.config.output.include_text,
                    include_embedding=self.config.output.include_embeddings,
                    include_metadata=self.config.output.include_metadata
                )
                for chunk in result.chunks
            ],
            "statistics": result.stats.to_dict() if self.config.output.include_stats else None,
            "processing_timestamp": result.processing_timestamp,
            "config_used": result.config_used if self.config.output.include_metadata else None
        }
        
        async with aiofiles.open(output_path, 'w', encoding='utf-8') as f:
            await f.write(json.dumps(output_data, indent=2))
    
    async def _write_parquet_output(self, result: ProcessingResult, output_path: Path) -> None:
        """Write results in Parquet format."""
        try:
            import pandas as pd
            import pyarrow as pa
            import pyarrow.parquet as pq
            
            # Convert chunks to pandas DataFrame
            chunk_data = []
            for chunk in result.chunks:
                chunk_dict = chunk.to_dict(
                    include_text=self.config.output.include_text,
                    include_embedding=self.config.output.include_embeddings,
                    include_metadata=self.config.output.include_metadata
                )
                chunk_data.append(chunk_dict)
            
            df = pd.DataFrame(chunk_data)
            
            # Write to Parquet
            table = pa.Table.from_pandas(df)
            pq.write_table(table, output_path)
            
        except ImportError:
            logger.error("Parquet dependencies (pandas, pyarrow) not available")
            # Fallback to JSON
            await self._write_json_output(result, output_path.with_suffix('.json'))
    
    def get_stats(self) -> ProcessingStats:
        """Get current processing statistics."""
        return self._stats
    
    def get_normalizer(self) -> TextNormalizer:
        """Get the text normalizer for direct access."""
        return self._normalizer
    
    def get_cache_info(self) -> Dict[str, Any]:
        """Get cache information and statistics."""
        return {
            "cached_chunks": len(self._content_cache),
            "cache_memory_usage": sum(
                len(chunk.text) + len(chunk.text_hash) 
                for chunk in self._content_cache.values()
            ),
            "embedding_model_loaded": self._embedding_model is not None,
        }
    
    def clear_cache(self) -> None:
        """Clear all caches to free memory."""
        self._content_cache.clear()
        
        # Clear embedding model cache
        if self._embedding_model:
            try:
                # Clear model cache if available
                if hasattr(self._embedding_model, 'clear_cache'):
                    self._embedding_model.clear_cache()
                
                # Clear tokenizer cache if available
                if hasattr(self._embedding_model, '_tokenizer'):
                    tokenizer = self._embedding_model._tokenizer
                    if hasattr(tokenizer, 'clean_up_tokenization_spaces'):
                        # Clear tokenizer caches
                        pass
                
                # Force garbage collection for model weights
                import gc
                import torch
                if torch.cuda.is_available():
                    torch.cuda.empty_cache()
                gc.collect()
                
                logger.info("Embedding model cache cleared")
            except Exception as e:
                logger.warning(f"Error clearing embedding model cache: {e}")
        
        # Clear sentence and chunk deduplicator caches
        if hasattr(self._sentence_dedup, 'clear_cache'):
            self._sentence_dedup.clear_cache()
        if hasattr(self._chunk_dedup, 'clear_cache'):
            self._chunk_dedup.clear_cache()
        
        logger.info("Pipeline caches cleared")


# Future enhancements planned for v2:
# - File watching for real-time processing with inotify/watchdog
# - URL processing support with aiohttp and content-type detection
# - Distributed processing with Celery/RQ for multiple workers
# - Progress callbacks for UI integration and WebSocket updates
# - Incremental processing for large datasets with checkpoints
# - Custom processing plugins with plugin architecture
# - Automatic retry logic with exponential backoff
# - Memory usage monitoring and optimization with profiling
# - Streaming API for very large documents
# - Parallel chunk processing with work stealing queues 